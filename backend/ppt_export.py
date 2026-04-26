import io
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm


NAVY   = RGBColor(0x1E, 0x3A, 0x8A)
BLUE   = RGBColor(0x3B, 0x82, 0xF6)
LIGHT  = RGBColor(0xEF, 0xF6, 0xFF)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
GREEN  = RGBColor(0x05, 0x96, 0x69)
DARK   = RGBColor(0x11, 0x18, 0x27)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)

W = Inches(13.33)
H = Inches(7.5)


def _slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        sp = ph._element
        sp.getparent().remove(sp)
    return slide


def _rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def _txbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(x, y, w, h)


def _para(tf, text, size=Pt(12), bold=False, color=DARK, align=PP_ALIGN.LEFT, space_before=Pt(0), italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    if text:
        run = p.add_run()
        run.text = text
        run.font.size = size
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return p


def _strip_md(text):
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    text = re.sub(r'#{1,3}\s*', '', text)
    text = re.sub(r'^[-•]\s*', '', text, flags=re.MULTILINE)
    return text.strip()


def _slide_title_page(prs, run_id):
    slide = _slide(prs)
    _rect(slide, 0, 0, W, H, fill=NAVY)
    _rect(slide, 0, Inches(5.8), W, Inches(1.7), fill=RGBColor(0x17, 0x2D, 0x70))

    accent = _txbox(slide, Inches(0.6), Inches(0.5), Inches(6), Inches(0.4))
    tf = accent.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "INDUSTRY INTELLIGENCE REPORT"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    r.font.italic = False

    title = _txbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(1.8))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "FMCG M&A Intelligence"
    r.font.size = Pt(48)
    r.font.bold = True
    r.font.color.rgb = WHITE

    sub = _txbox(slide, Inches(0.6), Inches(3.0), Inches(10), Inches(0.6))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Newsletter"
    r.font.size = Pt(36)
    r.font.bold = False
    r.font.color.rgb = BLUE

    date_box = _txbox(slide, Inches(0.6), Inches(4.0), Inches(8), Inches(0.4))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"Period: {datetime.now().strftime('%B %Y')}  ·  Powered by LangGraph + Claude Sonnet"
    r.font.size = Pt(13)
    r.font.color.rgb = RGBColor(0xAB, 0xBD, 0xE3)

    divider = _txbox(slide, Inches(0.6), Inches(4.6), Inches(4), Inches(0.05))
    tf = divider.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "─" * 28
    r.font.color.rgb = BLUE
    r.font.size = Pt(10)

    footer = _txbox(slide, Inches(0.6), Inches(6.1), Inches(10), Inches(0.35))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"Run ID: {run_id[:16]}  ·  Confidential — For Internal Use"
    r.font.size = Pt(9)
    r.font.color.rgb = RGBColor(0x8B, 0x9B, 0xC0)


def _slide_executive_summary(prs, sections):
    slide = _slide(prs)
    _rect(slide, 0, 0, W, Inches(1.3), fill=NAVY)
    _rect(slide, 0, Inches(1.3), W, H - Inches(1.3), fill=WHITE)

    label = _txbox(slide, Inches(0.5), Inches(0.18), Inches(4), Inches(0.35))
    tf = label.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "01  /  EXECUTIVE SUMMARY"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = ACCENT

    heading = _txbox(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.65))
    tf = heading.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Market Snapshot"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE

    summary = _strip_md(sections.get("executive_summary", "No summary available."))
    sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', summary) if s.strip()]

    y = Inches(1.55)
    for sent in sentences[:5]:
        _rect(slide, Inches(0.5), y, Inches(0.06), Inches(0.28), fill=BLUE)
        tb = _txbox(slide, Inches(0.75), y - Inches(0.04), Inches(11.8), Inches(0.45))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = sent
        r.font.size = Pt(14)
        r.font.color.rgb = DARK
        y += Inches(0.62)

    date_tag = _txbox(slide, Inches(10.5), Inches(6.9), Inches(2.5), Inches(0.35))
    tf = date_tag.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = datetime.now().strftime("%B %Y")
    r.font.size = Pt(9)
    r.font.color.rgb = GRAY


def _slide_top_deals(prs, sections):
    slide = _slide(prs)
    _rect(slide, 0, 0, W, Inches(1.3), fill=NAVY)
    _rect(slide, 0, Inches(1.3), W, H - Inches(1.3), fill=LIGHT)

    label = _txbox(slide, Inches(0.5), Inches(0.18), Inches(4), Inches(0.35))
    tf = label.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "02  /  TOP DEALS"
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = ACCENT

    heading = _txbox(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.65))
    tf = heading.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Key Transactions This Period"
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE

    deals = sections.get("top_deals", [])
    if isinstance(deals, str):
        deals = [d.strip() for d in deals.split("\n") if d.strip()]

    col_w = Inches(4.0)
    col_h = Inches(2.3)
    positions = [
        (Inches(0.4), Inches(1.45)),
        (Inches(4.65), Inches(1.45)),
        (Inches(8.9), Inches(1.45)),
        (Inches(0.4), Inches(3.95)),
        (Inches(4.65), Inches(3.95)),
        (Inches(8.9), Inches(3.95)),
    ]

    for i, deal in enumerate(deals[:6]):
        if i >= len(positions):
            break
        x, y = positions[i]
        _rect(slide, x, y, col_w, col_h, fill=WHITE,
              line=RGBColor(0xD1, 0xD5, 0xDB), line_w=Pt(0.5))
        _rect(slide, x, y, col_w, Inches(0.07), fill=BLUE)

        text = _strip_md(str(deal))[:180]
        num = _txbox(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.35))
        tf = num.text_frame
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = f"{i+1:02d}"
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = BLUE

        tb = _txbox(slide, x + Inches(0.15), y + Inches(0.5), col_w - Inches(0.3), col_h - Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = text
        r.font.size = Pt(10.5); r.font.color.rgb = DARK


def _slide_sector_breakdown(prs, sections):
    slide = _slide(prs)
    _rect(slide, 0, 0, W, Inches(1.3), fill=NAVY)
    _rect(slide, 0, Inches(1.3), W, H - Inches(1.3), fill=WHITE)

    label = _txbox(slide, Inches(0.5), Inches(0.18), Inches(5), Inches(0.35))
    tf = label.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "03  /  SECTOR BREAKDOWN"
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = ACCENT

    heading = _txbox(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.65))
    tf = heading.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Activity by Category"
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE

    sectors_raw = sections.get("sector_breakdown", {})
    if isinstance(sectors_raw, dict):
        sectors = list(sectors_raw.items())[:4]
    else:
        text = _strip_md(str(sectors_raw))
        sectors = [("Sector Overview", text)]

    sector_colors = [NAVY, BLUE, GREEN, RGBColor(0x7C, 0x3A, 0xED)]
    col_w = Inches(2.9)
    col_h = Inches(4.8)
    gap = Inches(0.35)

    for i, (sector, content) in enumerate(sectors):
        x = Inches(0.4) + i * (col_w + gap)
        y = Inches(1.45)
        color = sector_colors[i % len(sector_colors)]

        _rect(slide, x, y, col_w, col_h, fill=LIGHT,
              line=RGBColor(0xE5, 0xE7, 0xEB), line_w=Pt(0.5))
        _rect(slide, x, y, col_w, Inches(0.55), fill=color)

        hd = _txbox(slide, x + Inches(0.12), y + Inches(0.08), col_w - Inches(0.24), Inches(0.4))
        tf = hd.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = _strip_md(str(sector))[:35]
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE

        body_text = _strip_md(str(content)) if content else "No deals reported this period."
        items = [s.strip() for s in re.split(r'[.\n]', body_text) if len(s.strip()) > 10]

        tb = _txbox(slide, x + Inches(0.12), y + Inches(0.7), col_w - Inches(0.24), col_h - Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        for item in items[:5]:
            p = tf.add_paragraph()
            p.space_before = Pt(4)
            r = p.add_run()
            r.text = f"• {item[:100]}"
            r.font.size = Pt(9.5); r.font.color.rgb = DARK


def _slide_deals_to_watch(prs, sections):
    slide = _slide(prs)
    _rect(slide, 0, 0, W, Inches(1.3), fill=NAVY)
    _rect(slide, 0, Inches(1.3), W, H - Inches(1.3), fill=WHITE)

    label = _txbox(slide, Inches(0.5), Inches(0.18), Inches(5), Inches(0.35))
    tf = label.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "04  /  DEALS TO WATCH"
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = ACCENT

    heading = _txbox(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.65))
    tf = heading.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Emerging Themes & Situations"
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE

    deals_raw = sections.get("deals_to_watch", [])
    if isinstance(deals_raw, list):
        items = [_strip_md(str(d)) for d in deals_raw if str(d).strip()]
    else:
        items = [s.strip() for s in str(deals_raw).split("\n") if s.strip()]

    icon_colors = [BLUE, GREEN, RGBColor(0xF5, 0x9E, 0x0B)]

    y = Inches(1.6)
    for i, item in enumerate(items[:5]):
        color = icon_colors[i % len(icon_colors)]
        _rect(slide, Inches(0.5), y + Inches(0.05), Inches(0.45), Inches(0.45), fill=color)

        num = _txbox(slide, Inches(0.5), y + Inches(0.05), Inches(0.45), Inches(0.45))
        tf = num.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(i + 1)
        r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE

        tb = _txbox(slide, Inches(1.1), y, Inches(11.6), Inches(0.65))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = item[:220]
        r.font.size = Pt(12.5); r.font.color.rgb = DARK
        y += Inches(0.95)


def _slide_market_outlook(prs, sections):
    slide = _slide(prs)
    _rect(slide, 0, 0, W, H, fill=NAVY)
    _rect(slide, 0, Inches(5.5), W, Inches(2.0), fill=RGBColor(0x17, 0x2D, 0x70))

    label = _txbox(slide, Inches(0.6), Inches(0.4), Inches(5), Inches(0.35))
    tf = label.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "05  /  MARKET OUTLOOK"
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = ACCENT

    heading = _txbox(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8))
    tf = heading.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "What to Expect"
    r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = WHITE

    outlook = _strip_md(sections.get("market_outlook", "Continued M&A activity expected across FMCG categories."))
    sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', outlook) if s.strip()]

    y = Inches(2.0)
    for sent in sentences[:4]:
        _rect(slide, Inches(0.6), y + Inches(0.1), Inches(0.08), Inches(0.3), fill=ACCENT)
        tb = _txbox(slide, Inches(0.9), y, Inches(11.8), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = sent
        r.font.size = Pt(15); r.font.color.rgb = RGBColor(0xD1, 0xE4, 0xFF)
        y += Inches(0.75)

    footer = _txbox(slide, Inches(0.6), Inches(6.15), Inches(10), Inches(0.35))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"FMCG M&A Intelligence · {datetime.now().strftime('%B %Y')} · Confidential"
    r.font.size = Pt(9); r.font.color.rgb = RGBColor(0x8B, 0x9B, 0xC0)


def _slide_pipeline(prs, pipeline_stats):
    slide = _slide(prs)
    _rect(slide, 0, 0, W, Inches(1.3), fill=NAVY)
    _rect(slide, 0, Inches(1.3), W, H - Inches(1.3), fill=WHITE)

    label = _txbox(slide, Inches(0.5), Inches(0.18), Inches(5), Inches(0.35))
    tf = label.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "APPENDIX  /  PIPELINE"
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = ACCENT

    heading = _txbox(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.65))
    tf = heading.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "How the Pipeline Works"
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE

    stages = [
        ("Ingestion", f"{pipeline_stats.get('raw', '~80')} articles", "Tavily Search · 11 dynamic queries · 30-day window"),
        ("Deduplication", "Jaccard >70%", "Title word-set overlap · near-duplicates removed"),
        ("Relevance Filter", "Score ≥ 0.15", "FMCG + deal keywords · weighted density scoring"),
        ("Credibility Check", f"{pipeline_stats.get('after_credibility', '~15')} pass", "33-domain trust map · social media blocked"),
        ("Enrichment", "LLM extract", "Claude: deal type · companies · deal value"),
        ("Newsletter", "5 sections", "Bloomberg-style brief · factual · named companies"),
    ]

    col_w = Inches(1.95)
    arrow_w = Inches(0.25)
    x = Inches(0.3)
    y_box = Inches(1.6)

    for i, (title, metric, desc) in enumerate(stages):
        _rect(slide, x, y_box, col_w, Inches(2.5), fill=LIGHT,
              line=RGBColor(0xD1, 0xD5, 0xDB), line_w=Pt(0.5))
        _rect(slide, x, y_box, col_w, Inches(0.55), fill=NAVY)

        hd = _txbox(slide, x + Inches(0.08), y_box + Inches(0.08), col_w - Inches(0.16), Inches(0.4))
        tf = hd.text_frame
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = f"{i+1}. {title}"
        r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = WHITE

        met = _txbox(slide, x + Inches(0.08), y_box + Inches(0.65), col_w - Inches(0.16), Inches(0.4))
        tf = met.text_frame
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = metric
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = BLUE

        db = _txbox(slide, x + Inches(0.08), y_box + Inches(1.15), col_w - Inches(0.16), Inches(1.2))
        tf = db.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = desc
        r.font.size = Pt(9); r.font.color.rgb = GRAY

        x += col_w
        if i < len(stages) - 1:
            arr = _txbox(slide, x, y_box + Inches(1.0), arrow_w, Inches(0.4))
            tf = arr.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = ">"
            r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = BLUE
            x += arrow_w

    cred_box = _txbox(slide, Inches(0.3), Inches(4.3), Inches(12.7), Inches(2.8))
    tf = cred_box.text_frame
    tf.word_wrap = True
    _para(tf, "Credibility Scoring — Transparent Assumptions", size=Pt(10), bold=True, color=NAVY, space_before=Pt(0))
    _para(tf, "Reuters/Bloomberg: 0.95  ·  WSJ/FT: 0.92  ·  CNBC: 0.90  ·  Forbes: 0.88  ·  Wire services: 0.85–0.88", size=Pt(9), color=GRAY)
    _para(tf, "FMCG trade press (fooddive, grocerygazette, just-food): 0.78–0.80  ·  Advisory sites: 0.65–0.76  ·  Unknown: 0.46  ·  Social media: 0.20", size=Pt(9), color=GRAY)
    _para(tf, "Threshold: score > 0.45 to pass. Social media always blocked regardless of relevance score.", size=Pt(9), color=GRAY, italic=True)


def generate_pptx(newsletter_draft: str, sections: dict, articles: list, pipeline_stats: dict, run_id: str) -> bytes:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    _slide_title_page(prs, run_id)
    _slide_executive_summary(prs, sections)
    _slide_top_deals(prs, sections)
    _slide_sector_breakdown(prs, sections)
    _slide_deals_to_watch(prs, sections)
    _slide_market_outlook(prs, sections)
    _slide_pipeline(prs, pipeline_stats)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
